import streamlit as st
import os
import pickle
import hashlib
import json
import time
from pathlib import Path
from dotenv import load_dotenv
import glob
import tempfile
import io
# LangChain imports
from langchain_groq import ChatGroq
from langchain_community.llms import Ollama
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_classic.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain_classic.chains.retrieval import create_retrieval_chain
from langchain_community.vectorstores import FAISS
from langchain_community.retrievers import BM25Retriever
from langchain_classic.retrievers import EnsembleRetriever
from langchain_core.documents import Document

# Multi-format document processing imports
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None

try:
    import whisper
except ImportError:
    whisper = None

try:
    from pydub import AudioSegment
except ImportError:
    AudioSegment = None

try:
    from moviepy.editor import VideoFileClip
except ImportError:
    VideoFileClip = None

# Load environment
load_dotenv()

# Configuration
EMBEDDING_MODEL = "all-MiniLM-L6-v2"
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
SEMANTIC_WEIGHT = 0.7
KEYWORD_WEIGHT = 0.3
NUM_RESULTS = 8

# Available Models
AVAILABLE_MODELS = {
    "Local Models (Ollama)": {
        "llama3.2:3b": "llama3.2:3b",
        "llama3.2:1b": "llama3.2:1b",
        "llama3.1:8b": "llama3.1:8b",
        "mistral:7b": "mistral:7b",
        "codellama:7b": "codellama:7b"
    },
    "API Models (Groq)": {
        "Llama 3.1 8B": "llama-3.1-8b-instant",
        "Llama 3.1 70B": "llama-3.1-70b-versatile",
        "Mixtral 8x7B": "mixtral-8x7b-32768",
        "Gemma 2 9B": "gemma2-9b-it",
        "GPT-OSS 120B": "openai/gpt-oss-120b"
    }
}

def check_ollama_connection():
    """Check if Ollama is running and accessible"""
    try:
        import requests
        response = requests.get("http://localhost:11434/api/tags", timeout=5)
        return response.status_code == 200
    except:
        return False

def get_available_ollama_models():
    """Get list of available Ollama models"""
    try:
        import requests
        response = requests.get("http://localhost:11434/api/tags", timeout=5)
        if response.status_code == 200:
            models_data = response.json()
            return [model['name'] for model in models_data.get('models', [])]
        return []
    except:
        return []

def initialize_llm(model_provider, model_name, groq_api_key=None):
    """Initialize the appropriate LLM based on provider"""
    try:
        if model_provider == "Local Models (Ollama)":
            if not check_ollama_connection():
                st.error("❌ Ollama is not running! Please start Ollama service.")
                return None

            available_models = get_available_ollama_models()
            if model_name not in available_models:
                st.warning(f"⚠️ Model '{model_name}' not found in Ollama.")
                st.info("To install: `ollama pull " + model_name + "`")
                return None

            return Ollama(model=model_name, base_url="http://localhost:11434", temperature=0.1)

        elif model_provider == "API Models (Groq)":
            if not groq_api_key:
                st.error("❌ GROQ_API_KEY not found!")
                return None

            return ChatGroq(groq_api_key=groq_api_key, model_name=model_name, temperature=0)

        return None
    except Exception as e:
        st.error(f"❌ Error initializing LLM: {e}")
        return None

groq_api_key = os.environ.get('GROQ_API_KEY')

def clean_text(text):
    """Clean and validate text content"""
    if not text or not isinstance(text, str):
        return ""
    cleaned = text.encode('utf-8', errors='ignore').decode('utf-8').strip()
    return ' '.join(cleaned.split())

# Supported file extensions
SUPPORTED_EXTENSIONS = {
    'pdf': ['.pdf'],
    'word': ['.docx', '.doc'],
    'excel': ['.xlsx', '.xls'],
    'powerpoint': ['.pptx', '.ppt'],
    'text': ['.txt', '.md', '.markdown'],
    'image': ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'],
    'audio': ['.mp3', '.wav', '.m4a', '.flac', '.ogg'],
    'video': ['.mp4', '.avi', '.mov', '.mkv', '.flv']
}

def get_all_supported_extensions():
    """Get flat list of all supported extensions"""
    return [ext for exts in SUPPORTED_EXTENSIONS.values() for ext in exts]

def get_all_documents(folder_path):
    """Get all supported documents from specified folder and subfolders"""
    all_files = []
    supported_exts = get_all_supported_extensions()
    
    for ext in supported_exts:
        patterns = [
            os.path.join(folder_path, f"**/*{ext}"),
            os.path.join(folder_path, f"*{ext}")
        ]
        for pattern in patterns:
            all_files.extend(glob.glob(pattern, recursive=True))
    
    return [Path(f) for f in set(all_files) if os.path.exists(f)]

def get_all_pdfs(folder_path):
    """Get all PDFs from specified folder and subfolders (backward compatibility)"""
    pdf_patterns = [
        os.path.join(folder_path, "**/*.pdf"),
        os.path.join(folder_path, "*.pdf")
    ]
    all_pdfs = []
    for pattern in pdf_patterns:
        all_pdfs.extend(glob.glob(pattern, recursive=True))
    return [Path(pdf) for pdf in set(all_pdfs) if os.path.exists(pdf)]

def load_pdf_document(file_path):
    """Load PDF document"""
    loader = PyPDFLoader(str(file_path))
    return loader.load()

def load_word_document(file_path):
    """Load Word document"""
    if DocxDocument is None:
        raise ImportError("python-docx not installed. Run: pip install python-docx")
    
    doc = DocxDocument(str(file_path))
    text_content = []
    
    for i, para in enumerate(doc.paragraphs):
        if para.text.strip():
            text_content.append(para.text)
    
    full_text = '\n\n'.join(text_content)
    return [Document(
        page_content=full_text,
        metadata={'source': str(file_path), 'page': 1}
    )]

def load_excel_document(file_path):
    """Load Excel document"""
    if load_workbook is None:
        raise ImportError("openpyxl not installed. Run: pip install openpyxl")
    
    wb = load_workbook(str(file_path), read_only=True)
    documents = []
    
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows_text = []
        
        for row in sheet.iter_rows(values_only=True):
            row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
            if row_text.strip():
                rows_text.append(row_text)
        
        if rows_text:
            sheet_content = '\n'.join(rows_text)
            documents.append(Document(
                page_content=sheet_content,
                metadata={'source': str(file_path), 'sheet': sheet_name, 'page': 1}
            ))
    
    return documents

def load_powerpoint_document(file_path):
    """Load PowerPoint document"""
    if Presentation is None:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")
    
    prs = Presentation(str(file_path))
    documents = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        
        if slide_text:
            documents.append(Document(
                page_content='\n'.join(slide_text),
                metadata={'source': str(file_path), 'slide': i + 1, 'page': i + 1}
            ))
    
    return documents

def load_text_document(file_path):
    """Load text or markdown document"""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        content = f.read()
    
    return [Document(
        page_content=content,
        metadata={'source': str(file_path), 'page': 1}
    )]

def load_image_document(file_path):
    """Load image and extract text using OCR"""
    if Image is None or pytesseract is None:
        raise ImportError("Pillow or pytesseract not installed. Run: pip install Pillow pytesseract")
    
    try:
        img = Image.open(str(file_path))
        text = pytesseract.image_to_string(img)
        
        return [Document(
            page_content=text if text.strip() else f"[Image: {file_path.name}]",
            metadata={'source': str(file_path), 'type': 'image', 'page': 1}
        )]
    except Exception as e:
        return [Document(
            page_content=f"[Image processing failed: {str(e)}]",
            metadata={'source': str(file_path), 'type': 'image', 'page': 1}
        )]

def load_audio_document(file_path):
    """Load audio and transcribe using Whisper"""
    if whisper is None:
        raise ImportError("openai-whisper not installed. Run: pip install openai-whisper")
    
    try:
        model = whisper.load_model("base")
        result = model.transcribe(str(file_path))
        
        return [Document(
            page_content=result["text"],
            metadata={'source': str(file_path), 'type': 'audio', 'page': 1}
        )]
    except Exception as e:
        return [Document(
            page_content=f"[Audio transcription failed: {str(e)}]",
            metadata={'source': str(file_path), 'type': 'audio', 'page': 1}
        )]

def load_video_document(file_path):
    """Load video and extract audio for transcription"""
    if whisper is None or VideoFileClip is None:
        raise ImportError("Required libraries not installed. Run: pip install openai-whisper moviepy")
    
    try:
        # Extract audio from video
        video = VideoFileClip(str(file_path))
        
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_audio:
            temp_audio_path = temp_audio.name
            video.audio.write_audiofile(temp_audio_path, verbose=False, logger=None)
        
        video.close()
        
        # Transcribe audio
        model = whisper.load_model("base")
        result = model.transcribe(temp_audio_path)
        
        # Clean up temp file
        os.unlink(temp_audio_path)
        
        return [Document(
            page_content=result["text"],
            metadata={'source': str(file_path), 'type': 'video', 'page': 1}
        )]
    except Exception as e:
        return [Document(
            page_content=f"[Video transcription failed: {str(e)}]",
            metadata={'source': str(file_path), 'type': 'video', 'page': 1}
        )]

def load_document_universal(file_path):
    """Universal document loader that detects file type and loads accordingly"""
    file_path = Path(file_path)
    extension = file_path.suffix.lower()
    
    try:
        if extension == '.pdf':
            return load_pdf_document(file_path)
        elif extension in ['.docx', '.doc']:
            return load_word_document(file_path)
        elif extension in ['.xlsx', '.xls']:
            return load_excel_document(file_path)
        elif extension in ['.pptx', '.ppt']:
            return load_powerpoint_document(file_path)
        elif extension in ['.txt', '.md', '.markdown']:
            return load_text_document(file_path)
        elif extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']:
            return load_image_document(file_path)
        elif extension in ['.mp3', '.wav', '.m4a', '.flac', '.ogg']:
            return load_audio_document(file_path)
        elif extension in ['.mp4', '.avi', '.mov', '.mkv', '.flv']:
            return load_video_document(file_path)
        else:
            # Try as text file
            return load_text_document(file_path)
    except Exception as e:
        st.warning(f"⚠️ Error loading {file_path.name}: {str(e)}")
        return []

def scan_folder_for_databases(folder_path):
    """Scan a folder for vector databases"""
    databases = []
    try:
        if not os.path.exists(folder_path):
            return databases

        vector_db_path = os.path.join(folder_path, "vector_databases")
        if os.path.exists(vector_db_path):
            for item in os.listdir(vector_db_path):
                if item.startswith('db_'):
                    db_path = os.path.join(vector_db_path, item)
                    metadata_path = os.path.join(db_path, 'metadata.json')
                    if os.path.exists(metadata_path):
                        try:
                            with open(metadata_path, 'r') as f:
                                metadata = json.load(f)
                            
                            # Support both old (total_pdfs) and new (total_files) format
                            total_files = metadata.get('total_files', metadata.get('total_pdfs', 0))
                            file_types = metadata.get('file_types', {})
                            
                            if file_types:
                                type_summary = ', '.join([f"{count} {ftype}" for ftype, count in file_types.items()])
                                name = f"{metadata.get('embedding_model', 'Unknown')} - {type_summary}"
                            else:
                                name = f"{metadata.get('embedding_model', 'Unknown')} - {total_files} files"
                            
                            databases.append({
                                'name': name,
                                'path': vector_db_path,
                                'metadata': metadata,
                                'db_folder': item
                            })
                        except:
                            continue

    except Exception:
        pass

    return databases

def load_existing_database(db_path, db_folder):
    """Load an existing vector database"""
    try:
        embeddings = HuggingFaceEmbeddings(
            model_name=f"sentence-transformers/{EMBEDDING_MODEL}",
            model_kwargs={'device': 'cpu'},
            encode_kwargs={'normalize_embeddings': True}
        )

        base_path = os.path.join(db_path, db_folder)
        faiss_path = os.path.join(base_path, "faiss_index")
        bm25_path = os.path.join(base_path, "bm25_retriever.pkl")
        documents_path = os.path.join(base_path, "documents.pkl")
        metadata_path = os.path.join(base_path, "metadata.json")

        vector_store = FAISS.load_local(faiss_path, embeddings, allow_dangerous_deserialization=True)

        with open(bm25_path, 'rb') as f:
            bm25_retriever = pickle.load(f)

        with open(documents_path, 'rb') as f:
            documents = pickle.load(f)

        with open(metadata_path, 'r') as f:
            metadata = json.load(f)

        return vector_store, bm25_retriever, documents, embeddings, metadata

    except Exception as e:
        st.error(f"Error loading database: {e}")
        return None

def create_new_database(file_list, save_folder):
    """Create new vector database from multiple file types"""
    try:
        embeddings = HuggingFaceEmbeddings(
            model_name=f"sentence-transformers/{EMBEDDING_MODEL}",
            model_kwargs={'device': 'cpu'},
            encode_kwargs={'normalize_embeddings': True}
        )

        with st.spinner(f"🔄 Processing {len(file_list)} files..."):
            documents = []
            progress_bar = st.progress(0)
            file_types = {}

            for i, file_path in enumerate(file_list):
                try:
                    file_ext = file_path.suffix.lower()
                    file_types[file_ext] = file_types.get(file_ext, 0) + 1
                    
                    st.text(f"📄 Loading {file_path.name} ({i+1}/{len(file_list)})...")
                    
                    # Use universal loader
                    loaded_docs = load_document_universal(file_path)

                    valid_docs = []
                    for doc in loaded_docs:
                        cleaned_content = clean_text(doc.page_content)
                        if len(cleaned_content) > 20:
                            doc.page_content = cleaned_content
                            if 'source_file' not in doc.metadata:
                                doc.metadata['source_file'] = file_path.name
                            if 'full_path' not in doc.metadata:
                                doc.metadata['full_path'] = str(file_path)
                            doc.metadata['file_type'] = file_ext
                            valid_docs.append(doc)

                    documents.extend(valid_docs)
                    progress_bar.progress((i + 1) / len(file_list))

                except Exception as e:
                    st.warning(f"⚠️ Failed to load {file_path.name}: {str(e)}")

            if not documents:
                st.error("❌ No valid documents were loaded!")
                return None

            st.text("✂️ Splitting documents...")
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=CHUNK_SIZE,
                chunk_overlap=CHUNK_OVERLAP
            )
            final_documents = text_splitter.split_documents(documents)

            st.text("🔢 Creating vector database...")
            vector_store = FAISS.from_documents(final_documents, embeddings)

            st.text("🔍 Creating keyword index...")
            bm25_retriever = BM25Retriever.from_documents(final_documents)
            bm25_retriever.k = NUM_RESULTS

            st.text("💾 Saving database...")
            config_str = f"{sorted([str(p) for p in file_list])}_{EMBEDDING_MODEL}_{CHUNK_SIZE}_{CHUNK_OVERLAP}"
            config_hash = hashlib.md5(config_str.encode()).hexdigest()[:12]

            os.makedirs(save_folder, exist_ok=True)
            db_folder = f"db_{config_hash}"
            db_path = os.path.join(save_folder, db_folder)
            os.makedirs(db_path, exist_ok=True)

            vector_store.save_local(os.path.join(db_path, "faiss_index"))

            with open(os.path.join(db_path, "bm25_retriever.pkl"), 'wb') as f:
                pickle.dump(bm25_retriever, f)

            with open(os.path.join(db_path, "documents.pkl"), 'wb') as f:
                pickle.dump(final_documents, f)

            metadata = {
                'created_at': time.strftime('%Y-%m-%d %H:%M:%S'),
                'updated_at': time.strftime('%Y-%m-%d %H:%M:%S'),
                'embedding_model': EMBEDDING_MODEL,
                'chunk_size': CHUNK_SIZE,
                'chunk_overlap': CHUNK_OVERLAP,
                'total_files': len(file_list),
                'total_chunks': len(final_documents),
                'file_types': file_types,
                'files': [str(p) for p in file_list],
                'config_hash': config_hash,
                # Backward compatibility
                'total_pdfs': file_types.get('.pdf', 0),
                'pdf_files': [str(p) for p in file_list if p.suffix.lower() == '.pdf']
            }

            with open(os.path.join(db_path, "metadata.json"), 'w') as f:
                json.dump(metadata, f, indent=2)

            progress_bar.progress(1.0)
            st.success(f"✅ Database created successfully!")

            return vector_store, bm25_retriever, final_documents, embeddings, metadata

    except Exception as e:
        st.error(f"❌ Error creating database: {e}")
        return None

def add_files_to_database(new_files, db_path, db_folder, existing_vector_store, existing_documents, embeddings):
    """Add new files to an existing database"""
    try:
        with st.spinner(f"🔄 Adding {len(new_files)} new files..."):
            new_documents = []
            progress_bar = st.progress(0)
            
            # Load existing metadata
            metadata_path = os.path.join(db_path, db_folder, "metadata.json")
            with open(metadata_path, 'r') as f:
                metadata = json.load(f)
            
            existing_files = set(metadata.get('files', metadata.get('pdf_files', [])))
            file_types = metadata.get('file_types', {})

            for i, file_path in enumerate(new_files):
                try:
                    # Skip if file already in database
                    if str(file_path) in existing_files:
                        st.info(f"⏭️ Skipping {file_path.name} (already in database)")
                        progress_bar.progress((i + 1) / len(new_files))
                        continue
                    
                    file_ext = file_path.suffix.lower()
                    file_types[file_ext] = file_types.get(file_ext, 0) + 1
                    
                    st.text(f"📄 Loading {file_path.name} ({i+1}/{len(new_files)})...")
                    
                    # Use universal loader
                    loaded_docs = load_document_universal(file_path)

                    valid_docs = []
                    for doc in loaded_docs:
                        cleaned_content = clean_text(doc.page_content)
                        if len(cleaned_content) > 20:
                            doc.page_content = cleaned_content
                            if 'source_file' not in doc.metadata:
                                doc.metadata['source_file'] = file_path.name
                            if 'full_path' not in doc.metadata:
                                doc.metadata['full_path'] = str(file_path)
                            doc.metadata['file_type'] = file_ext
                            valid_docs.append(doc)

                    new_documents.extend(valid_docs)
                    existing_files.add(str(file_path))
                    progress_bar.progress((i + 1) / len(new_files))

                except Exception as e:
                    st.warning(f"⚠️ Failed to load {file_path.name}: {str(e)}")

            if not new_documents:
                st.warning("⚠️ No new documents to add!")
                return None

            st.text("✂️ Splitting new documents...")
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=CHUNK_SIZE,
                chunk_overlap=CHUNK_OVERLAP
            )
            final_new_documents = text_splitter.split_documents(new_documents)

            st.text("🔢 Updating vector database...")
            # Add new documents to existing vector store
            existing_vector_store.add_documents(final_new_documents)

            st.text("🔍 Updating keyword index...")
            # Combine all documents
            all_documents = existing_documents + final_new_documents
            bm25_retriever = BM25Retriever.from_documents(all_documents)
            bm25_retriever.k = NUM_RESULTS

            st.text("💾 Saving updated database...")
            base_path = os.path.join(db_path, db_folder)
            
            existing_vector_store.save_local(os.path.join(base_path, "faiss_index"))

            with open(os.path.join(base_path, "bm25_retriever.pkl"), 'wb') as f:
                pickle.dump(bm25_retriever, f)

            with open(os.path.join(base_path, "documents.pkl"), 'wb') as f:
                pickle.dump(all_documents, f)

            # Update metadata
            metadata.update({
                'updated_at': time.strftime('%Y-%m-%d %H:%M:%S'),
                'total_files': len(existing_files),
                'total_chunks': len(all_documents),
                'file_types': file_types,
                'files': list(existing_files),
                # Backward compatibility
                'total_pdfs': file_types.get('.pdf', 0),
                'pdf_files': [f for f in existing_files if f.endswith('.pdf')]
            })

            with open(metadata_path, 'w') as f:
                json.dump(metadata, f, indent=2)

            progress_bar.progress(1.0)
            st.success(f"✅ Added {len(final_new_documents)} new chunks from {len(new_documents)} documents!")

            return existing_vector_store, bm25_retriever, all_documents, embeddings, metadata

    except Exception as e:
        st.error(f"❌ Error adding files to database: {e}")
        return None

def process_answer_with_citations(answer, context_docs):
    """Add inline citations to the answer text"""
    import re

    citations = {}
    for i, doc in enumerate(context_docs):
        source_file = doc.metadata.get('source_file', 'Unknown')
        page_num = doc.metadata.get('page', 'Unknown')
        full_path = doc.metadata.get('full_path', '')

        citation_key = f"[{i+1}]"
        if full_path and os.path.exists(full_path):
            citation_link = f"""<sup><a href="file:///{full_path}" target="_blank" class="citation-link" title="Open {source_file} - Page {page_num}">[{i+1}]</a></sup>"""
        else:
            citation_link = f"""<sup class="citation-unavailable" title="{source_file} - Page {page_num}">[{i+1}]</sup>"""

        citations[citation_key] = citation_link

    processed_answer = answer

    patterns_to_cite = [
        r'(Item[‑\-]\s*No\.\s*[\d\-‑]+[\w\-‑]*)',
        r'(Rs\.?\s*[\d,]+\.?\d*)',
        r'(Chapter\s*[\d]+)',
        r'(\d+\.?\d*\s*%)',
        r'(Rule\s*[\d\.]+)',
        r'(\d+″\s*×\s*\d+″)',
        r'(Sq\s*[mft]+)',
        r'(MRS|CSR|Market[‑\-]Rate\s*System)',
    ]

    citation_usage = [False] * len(context_docs)

    for pattern in patterns_to_cite:
        matches = list(re.finditer(pattern, processed_answer, re.IGNORECASE))
        for match in reversed(matches):
            match_text = match.group(1).lower()
            best_citation_idx = None
            best_score = 0

            for i, doc in enumerate(context_docs):
                if citation_usage[i]:
                    continue

                doc_content = doc.page_content.lower()
                score = 0
                if 'item' in match_text and 'item' in doc_content:
                    score += 3
                if 'chapter' in match_text and 'chapter' in doc_content:
                    score += 3
                if 'rs' in match_text and 'rs' in doc_content:
                    score += 2
                if any(word in doc_content for word in match_text.split()):
                    score += 1

                if score > best_score:
                    best_score = score
                    best_citation_idx = i

            if best_citation_idx is not None and best_score > 0:
                citation_link = citations[f"[{best_citation_idx+1}]"]
                citation_usage[best_citation_idx] = True

                end_pos = match.end()
                processed_answer = (processed_answer[:end_pos] +
                                  citation_link +
                                  processed_answer[end_pos:])

    return processed_answer

def create_hybrid_retriever(vector_store, bm25_retriever):
    """Create hybrid retriever combining semantic and keyword search"""
    vector_retriever = vector_store.as_retriever(search_kwargs={"k": NUM_RESULTS})

    ensemble_retriever = EnsembleRetriever(
        retrievers=[vector_retriever, bm25_retriever],
        weights=[SEMANTIC_WEIGHT, KEYWORD_WEIGHT]
    )

    return ensemble_retriever

# Page Configuration
st.set_page_config(
    page_title="RAG Intelligence Hub",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# Initialize theme in session state
if 'app_theme' not in st.session_state:
    st.session_state.app_theme = 'dark'

# Theme selector in sidebar
with st.sidebar:
    st.markdown("### 🎨 Theme")
    theme_choice = st.radio(
        "Select Theme",
        ["🌙 Dark", "☀️ Light"],
        index=0 if st.session_state.app_theme == 'dark' else 1,
        key="theme_selector"
    )
    st.session_state.app_theme = 'dark' if theme_choice == "🌙 Dark" else 'light'

# Ultra Modern CSS with Complete Redesign - Dark/Light Theme Support
theme_class = "theme-light" if st.session_state.app_theme == 'light' else "theme-dark"

# Apply theme via JavaScript
st.markdown(f"""
<script>
    // Apply theme class to root element
    const root = document.documentElement;
    root.classList.remove('theme-dark', 'theme-light');
    root.classList.add('{theme_class}');
</script>
""", unsafe_allow_html=True)

# CSS Styles
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;500;600;700;800&family=JetBrains+Mono:wght@400;500;600&display=swap');

    * {
        font-family: 'Poppins', sans-serif;
        transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
    }

    /* Light Theme Background */
    html.theme-light .stApp {
        background: linear-gradient(135deg, #f8fafc 0%, #e2e8f0 50%, #cbd5e1 100%) !important;
    }
    
    /* Dark Theme Background (Default) */
    html.theme-dark .stApp,
    .stApp {
        background: linear-gradient(135deg, #0f172a 0%, #1e293b 50%, #334155 100%) !important;
    }

    .main {
        background: transparent;
        padding: 0 !important;
    }

    .main > div {
        background: transparent;
        padding: 0 !important;
    }

    /* Modern Header Bar */
    html.theme-dark .top-bar,
    .top-bar {
        background: rgba(255, 255, 255, 0.08);
        backdrop-filter: blur(20px);
        border-bottom: 1px solid rgba(99, 102, 241, 0.2);
        padding: 1.5rem 3rem;
        display: flex;
        justify-content: space-between;
        align-items: center;
        margin-bottom: 2rem;
    }
    
    html.theme-light .top-bar {
        background: rgba(255, 255, 255, 0.95);
        border-bottom: 1px solid rgba(99, 102, 241, 0.3);
    }

    .logo-section {
        display: flex;
        align-items: center;
        gap: 1rem;
    }

    .logo-icon {
        font-size: 2.5rem;
        background: linear-gradient(135deg, #6366f1, #a855f7);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        background-clip: text;
    }

    html.theme-dark .logo-text,
    .logo-text {
        color: white;
        font-size: 1.75rem;
        font-weight: 700;
        letter-spacing: -0.5px;
    }
    
    html.theme-light .logo-text {
        color: #1e293b;
    }

    .status-badge {
        background: linear-gradient(135deg, #10b981, #059669);
        color: white;
        padding: 0.5rem 1.25rem;
        border-radius: 50px;
        font-size: 0.875rem;
        font-weight: 600;
        display: inline-flex;
        align-items: center;
        gap: 0.5rem;
    }

    /* Container */
    .content-container {
        max-width: 1400px;
        margin: 0 auto;
        padding: 0 2rem 3rem 2rem;
    }

    /* Modern Cards */
    html.theme-dark .modern-card,
    .modern-card {
        background: rgba(255, 255, 255, 0.08);
        backdrop-filter: blur(20px);
        border: 1px solid rgba(99, 102, 241, 0.2);
        border-radius: 24px;
        padding: 2.5rem;
        margin-bottom: 2rem;
        box-shadow: 0 20px 60px rgba(0, 0, 0, 0.3);
    }
    
    html.theme-light .modern-card {
        background: rgba(255, 255, 255, 0.95);
        border: 1px solid rgba(99, 102, 241, 0.3);
        box-shadow: 0 8px 30px rgba(99, 102, 241, 0.15);
    }

    .modern-card:hover {
        transform: translateY(-4px);
        border-color: rgba(99, 102, 241, 0.6);
        box-shadow: 0 30px 80px rgba(99, 102, 241, 0.25);
    }

    html.theme-dark .card-title,
    .card-title {
        color: white;
        font-size: 1.5rem;
        font-weight: 700;
        margin-bottom: 1.5rem;
        display: flex;
        align-items: center;
        gap: 0.75rem;
    }
    
    html.theme-light .card-title {
        color: #1e293b;
    }

    .card-title-icon {
        font-size: 1.75rem;
    }

    /* Glassmorphism Panels */
    html.theme-dark .glass-panel,
    .glass-panel {
        background: rgba(255, 255, 255, 0.05);
        backdrop-filter: blur(10px);
        border: 1px solid rgba(99, 102, 241, 0.2);
        border-radius: 16px;
        padding: 1.5rem;
        margin-bottom: 1.5rem;
    }
    
    html.theme-light .glass-panel {
        background: rgba(255, 255, 255, 0.9);
        border: 1px solid rgba(99, 102, 241, 0.25);
    }

    /* Modern Tabs */
    .stTabs [data-baseweb="tab-list"] {
        background: rgba(255, 255, 255, 0.05);
        backdrop-filter: blur(10px);
        border-radius: 16px;
        padding: 0.5rem;
        gap: 0.5rem;
        border: 1px solid rgba(255, 255, 255, 0.1);
    }

    .stTabs [data-baseweb="tab"] {
        background: transparent;
        border-radius: 12px;
        color: rgba(255, 255, 255, 0.6) !important;
        font-weight: 600;
        padding: 0.875rem 1.5rem;
        border: none;
    }

    .stTabs [data-baseweb="tab"]:hover {
        background: rgba(255, 255, 255, 0.05);
        color: rgba(255, 255, 255, 0.9) !important;
    }

    .stTabs [aria-selected="true"] {
        background: linear-gradient(135deg, #6366f1, #a855f7) !important;
        color: white !important;
    }

    /* Text Colors - Adaptive */
    html.theme-dark p,
    html.theme-dark span:not(.status-badge span),
    html.theme-dark div:not(.logo-icon):not(.status-badge),
    html.theme-dark label,
    p, span, div, label {
        color: rgba(255, 255, 255, 0.9) !important;
    }
    
    html.theme-light p,
    html.theme-light span:not(.status-badge span),
    html.theme-light div:not(.logo-icon):not(.status-badge),
    html.theme-light label {
        color: rgba(30, 41, 59, 0.95) !important;
    }

    html.theme-dark h1,
    html.theme-dark h2,
    html.theme-dark h3,
    html.theme-dark h4,
    html.theme-dark h5,
    html.theme-dark h6,
    h1, h2, h3, h4, h5, h6 {
        color: white !important;
        font-weight: 700;
    }
    
    html.theme-light h1,
    html.theme-light h2,
    html.theme-light h3,
    html.theme-light h4,
    html.theme-light h5,
    html.theme-light h6 {
        color: #1e293b !important;
        font-weight: 700;
    }

    /* Modern Buttons */
    .stButton > button {
        background: linear-gradient(135deg, #6366f1, #a855f7);
        color: white !important;
        border: none;
        border-radius: 12px;
        padding: 0.875rem 2rem;
        font-weight: 600;
        font-size: 0.95rem;
        box-shadow: 0 8px 20px rgba(99, 102, 241, 0.4);
        letter-spacing: 0.3px;
    }

    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 12px 30px rgba(99, 102, 241, 0.5);
        background: linear-gradient(135deg, #a855f7, #6366f1);
    }

    /* Input Fields */
    .theme-dark .stTextInput > div > div > input,
    .theme-dark .stTextArea > div > div > textarea,
    .stTextInput > div > div > input,
    .stTextArea > div > div > textarea {
        background: rgba(255, 255, 255, 0.08) !important;
        border: 1px solid rgba(99, 102, 241, 0.3) !important;
        border-radius: 12px !important;
        color: white !important;
        padding: 0.875rem 1rem !important;
    }
    
    .theme-light .stTextInput > div > div > input,
    .theme-light .stTextArea > div > div > textarea {
        background: white !important;
        border: 1px solid rgba(99, 102, 241, 0.4) !important;
        color: #1e293b !important;
    }

    .stTextInput > div > div > input:focus,
    .stTextArea > div > div > textarea:focus {
        border-color: #6366f1 !important;
        box-shadow: 0 0 0 3px rgba(99, 102, 241, 0.2) !important;
    }

    .stTextInput > div > div > input::placeholder,
    .stTextArea > div > div > textarea::placeholder {
        color: rgba(99, 102, 241, 0.5) !important;
    }

    /* Select & Radio */
    .theme-dark .stSelectbox > div > div,
    .theme-dark .stRadio > div,
    .stSelectbox > div > div,
    .stRadio > div {
        background: rgba(255, 255, 255, 0.08) !important;
        border: 1px solid rgba(99, 102, 241, 0.3) !important;
        border-radius: 12px !important;
        color: white !important;
    }
    
    .theme-light .stSelectbox > div > div,
    .theme-light .stRadio > div {
        background: white !important;
        color: #1e293b !important;
    }

    .theme-dark .stSelectbox label,
    .theme-dark .stRadio label,
    .stSelectbox label,
    .stRadio label {
        color: rgba(255, 255, 255, 0.9) !important;
        font-weight: 500;
    }
    
    .theme-light .stSelectbox label,
    .theme-light .stRadio label {
        color: #1e293b !important;
    }

    /* Metrics */
    [data-testid="metric-container"] {
        background: rgba(255, 255, 255, 0.05);
        backdrop-filter: blur(10px);
        border: 1px solid rgba(255, 255, 255, 0.1);
        border-radius: 16px;
        padding: 1.5rem;
        box-shadow: 0 4px 12px rgba(0, 0, 0, 0.2);
    }

    [data-testid="metric-container"]:hover {
        transform: translateY(-4px);
        border-color: rgba(99, 102, 241, 0.5);
    }

    [data-testid="metric-container"] label {
        color: rgba(255, 255, 255, 0.6) !important;
        font-weight: 600;
        font-size: 0.75rem;
        text-transform: uppercase;
        letter-spacing: 1px;
    }

    [data-testid="metric-container"] [data-testid="stMetricValue"] {
        color: #6366f1 !important;
        font-weight: 800;
        font-size: 2rem;
    }

    /* Chat Messages */
    .theme-dark .stChatMessage,
    .stChatMessage {
        background: rgba(255, 255, 255, 0.08);
        backdrop-filter: blur(10px);
        border: 1px solid rgba(99, 102, 241, 0.2);
        border-radius: 16px;
        padding: 1.5rem;
        margin: 1rem 0;
    }
    
    .theme-light .stChatMessage {
        background: rgba(255, 255, 255, 0.95);
        border: 1px solid rgba(99, 102, 241, 0.3);
    }

    .theme-dark .stChatMessage p,
    .theme-dark .stChatMessage div,
    .theme-dark .stChatMessage span,
    .stChatMessage p,
    .stChatMessage div,
    .stChatMessage span {
        color: rgba(255, 255, 255, 0.9) !important;
    }
    
    .theme-light .stChatMessage p,
    .theme-light .stChatMessage div,
    .theme-light .stChatMessage span {
        color: rgba(30, 41, 59, 0.95) !important;
    }

    /* Chat Input */
    .theme-dark [data-testid="stChatInput"],
    [data-testid="stChatInput"] {
        background: rgba(255, 255, 255, 0.08) !important;
        backdrop-filter: blur(10px);
        border: 1px solid rgba(99, 102, 241, 0.3) !important;
        border-radius: 16px !important;
    }
    
    .theme-light [data-testid="stChatInput"] {
        background: white !important;
        border: 1px solid rgba(99, 102, 241, 0.4) !important;
    }

    .theme-dark [data-testid="stChatInput"] input,
    [data-testid="stChatInput"] input {
        background: transparent !important;
        color: white !important;
        border: none !important;
    }
    
    .theme-light [data-testid="stChatInput"] input {
        color: #1e293b !important;
    }

    /* Success/Info/Warning Boxes */
    .stSuccess, .stInfo, .stWarning, .stError {
        background: rgba(255, 255, 255, 0.05) !important;
        backdrop-filter: blur(10px);
        border-radius: 12px;
        border-left: 4px solid;
    }

    .stSuccess {
        border-color: #10b981;
    }

    .stInfo {
        border-color: #6366f1;
    }

    .stWarning {
        border-color: #f59e0b;
    }

    .stError {
        border-color: #ef4444;
    }

    /* Expander */
    .streamlit-expanderHeader {
        background: rgba(255, 255, 255, 0.05) !important;
        backdrop-filter: blur(10px);
        border-radius: 12px !important;
        color: white !important;
        font-weight: 600;
    }

    .streamlit-expanderContent {
        background: rgba(255, 255, 255, 0.03) !important;
        border-radius: 0 0 12px 12px !important;
    }

    /* Code Blocks */
    code, pre {
        background: rgba(0, 0, 0, 0.3) !important;
        color: #a855f7 !important;
        border-radius: 8px;
        padding: 0.25rem 0.5rem;
        font-family: 'JetBrains Mono', monospace;
    }

    /* Progress Bar */
    .stProgress > div > div {
        background: linear-gradient(90deg, #6366f1, #a855f7);
    }

    /* Citations */
    .citation-link {
        background: rgba(99, 102, 241, 0.2);
        color: #a78bfa !important;
        padding: 2px 8px;
        border-radius: 6px;
        text-decoration: none !important;
        font-weight: 600;
        border: 1px solid rgba(99, 102, 241, 0.3);
    }

    .citation-link:hover {
        background: linear-gradient(135deg, #6366f1, #a855f7);
        color: white !important;
        transform: translateY(-1px);
    }

    /* Scrollbar */
    ::-webkit-scrollbar {
        width: 10px;
        height: 10px;
    }

    ::-webkit-scrollbar-track {
        background: rgba(255, 255, 255, 0.05);
    }

    ::-webkit-scrollbar-thumb {
        background: linear-gradient(135deg, #6366f1, #a855f7);
        border-radius: 5px;
    }

    ::-webkit-scrollbar-thumb:hover {
        background: linear-gradient(135deg, #a855f7, #6366f1);
    }

    /* Answer Content */
    .theme-dark .answer-content,
    .answer-content {
        color: rgba(255, 255, 255, 0.9) !important;
        line-height: 1.8;
    }
    
    .theme-light .answer-content {
        color: #1e293b !important;
    }

    .theme-dark .answer-content h1,
    .theme-dark .answer-content h2,
    .theme-dark .answer-content h3,
    .answer-content h1,
    .answer-content h2,
    .answer-content h3 {
        color: #a78bfa !important;
        margin-top: 2rem;
        margin-bottom: 1rem;
    }
    
    .theme-light .answer-content h1,
    .theme-light .answer-content h2,
    .theme-light .answer-content h3 {
        color: #6366f1 !important;
    }

    .answer-content table {
        width: 100%;
        border-collapse: collapse;
        margin: 1.5rem 0;
    }

    .theme-dark .answer-content th,
    .theme-dark .answer-content td,
    .answer-content th,
    .answer-content td {
        border: 1px solid rgba(255, 255, 255, 0.1);
        padding: 0.875rem 1rem;
        text-align: left;
        color: rgba(255, 255, 255, 0.9) !important;
    }
    
    .theme-light .answer-content th,
    .theme-light .answer-content td {
        border: 1px solid rgba(99, 102, 241, 0.3);
        color: #1e293b !important;
    }

    .theme-dark .answer-content th,
    .answer-content th {
        background: rgba(99, 102, 241, 0.2);
        font-weight: 600;
        color: white !important;
    }
    
    .theme-light .answer-content th {
        color: #1e293b !important;
    }

    .answer-content ul,
    .answer-content ol {
        margin: 1rem 0;
        padding-left: 2rem;
    }

    .theme-dark .answer-content li,
    .answer-content li {
        margin: 0.5rem 0;
        color: rgba(255, 255, 255, 0.9) !important;
    }
    
    .theme-light .answer-content li {
        color: #1e293b !important;
    }

    /* Hide Streamlit Branding */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
</style>
""", unsafe_allow_html=True)

# Clear Streamlit cache on startup (force fresh scan)
st.cache_data.clear()
st.cache_resource.clear()

# Initialize Session State
if 'database_loaded' not in st.session_state:
    st.session_state.database_loaded = False
if 'llm_configured' not in st.session_state:
    st.session_state.llm_configured = False
if "messages" not in st.session_state:
    st.session_state.messages = []

# Modern Header Bar
st.markdown("""
<div class="top-bar">
    <div class="logo-section">
        <div class="logo-icon">🚀</div>
        <div class="logo-text">RAG Intelligence Hub</div>
    </div>
    <div class="status-badge">
        <span>●</span> System Ready
    </div>
</div>
""", unsafe_allow_html=True)

# Content Container
st.markdown('<div class="content-container">', unsafe_allow_html=True)

# Main Tabs Navigation
tab1, tab2, tab3 = st.tabs(["🏠 Dashboard", "⚙️ Configuration", "💬 Chat"])

# TAB 1: DASHBOARD
with tab1:
    st.markdown('<div class="card-title"><span class="card-title-icon">📊</span>System Overview</div>', unsafe_allow_html=True)

    if st.session_state.database_loaded and st.session_state.llm_configured:
        metadata = st.session_state.metadata

        col1, col2, col3, col4 = st.columns(4)

        with col1:
            st.metric("📄 Documents", metadata['total_pdfs'])
        with col2:
            st.metric("📝 Text Chunks", f"{metadata['total_chunks']:,}")
        with col3:
            st.metric("🤖 AI Model", st.session_state.selected_model.split(':')[0][:15])
        with col4:
            st.metric("🔍 Search Mode", "Hybrid")

        st.markdown('<div style="margin-top: 2rem;"></div>', unsafe_allow_html=True)

        col1, col2 = st.columns(2)

        with col1:
            st.markdown("""
            <div class="glass-panel">
                <h4 style="margin-bottom: 1rem;">📊 Database Details</h4>
                <p style="margin: 0.5rem 0; opacity: 0.85;">
                    <strong>Embedding Model:</strong> {}</p>
                <p style="margin: 0.5rem 0; opacity: 0.85;">
                    <strong>Created:</strong> {}</p>
                <p style="margin: 0.5rem 0; opacity: 0.85;">
                    <strong>Chunk Size:</strong> {} tokens</p>
            </div>
            """.format(
                metadata['embedding_model'],
                metadata['created_at'],
                metadata['chunk_size']
            ), unsafe_allow_html=True)

        with col2:
            st.markdown("""
            <div class="glass-panel">
                <h4 style="margin-bottom: 1rem;">🤖 AI Configuration</h4>
                <p style="margin: 0.5rem 0;">
                    <strong>Provider:</strong> {}</p>
                <p style="margin: 0.5rem 0;">
                    <strong>Model:</strong> {}</p>
                <p style="margin: 0.5rem 0;">
                    <strong>Status:</strong> <span style="color: #10b981;">● Active</span></p>
            </div>
            """.format(
                "Local (Ollama)" if st.session_state.model_provider == "Local Models (Ollama)" else "Cloud (Groq)",
                st.session_state.selected_model
            ), unsafe_allow_html=True)

    else:
        st.markdown("""
        <div class="modern-card">
            <h3 style="text-align: center; margin-bottom: 1rem;">Welcome to RAG Intelligence Hub</h3>
            <p style="text-align: center; font-size: 1.1rem; opacity: 0.8;">
                Please configure your database and AI model in the Configuration tab to get started.
            </p>
        </div>
        """, unsafe_allow_html=True)

# TAB 2: CONFIGURATION
with tab2:
    config_tab1, config_tab2 = st.tabs(["📂 Database", "🤖 AI Model"])

    with config_tab1:
        st.markdown('<div class="card-title"><span class="card-title-icon">📂</span>Database Management</div>', unsafe_allow_html=True)

        # Initialize session state for databases
        if 'scanned_databases' not in st.session_state:
            st.session_state.scanned_databases = []

        db_option = st.radio(
            "Choose Database Option",
            ["Load Existing Database", "Create New Database"],
            horizontal=True
        )

        if db_option == "Load Existing Database":
            st.markdown('<div style="margin-top: 2rem;"></div>', unsafe_allow_html=True)

            folder_path = st.text_input("📁 Database Folder Path", value="vector_databases")

            if st.button("🔍 Scan for Databases", key="scan_db"):
                databases = scan_folder_for_databases(".")
                st.session_state.scanned_databases = databases
                if databases:
                    st.success(f"✅ Found {len(databases)} database(s)")
                else:
                    st.warning("⚠️ No databases found")

            # Show database selection if databases are scanned
            if st.session_state.scanned_databases:
                st.markdown('<div style="margin-top: 1.5rem;"></div>', unsafe_allow_html=True)

                selected_db = st.selectbox(
                    "Select Database",
                    range(len(st.session_state.scanned_databases)),
                    format_func=lambda x: st.session_state.scanned_databases[x]['name']
                )

                col1, col2 = st.columns(2)
                with col1:
                    if st.button("🚀 Load Database", key="load_db", use_container_width=True):
                        with st.spinner("Loading database..."):
                            db_info = st.session_state.scanned_databases[selected_db]
                            result = load_existing_database(db_info['path'], db_info['db_folder'])

                            if result:
                                vector_store, bm25_retriever, documents, embeddings, metadata = result
                                st.session_state.vector_store = vector_store
                                st.session_state.bm25_retriever = bm25_retriever
                                st.session_state.documents = documents
                                st.session_state.embeddings = embeddings
                                st.session_state.metadata = metadata
                                st.session_state.database_loaded = True
                                st.session_state.current_db_info = db_info
                                st.success("✅ Database loaded successfully! Switch to AI Model tab to configure.")
                                st.balloons()
                
                with col2:
                    if st.button("➕ Add Files to Database", key="add_to_db_btn", use_container_width=True):
                        st.session_state.show_add_files = True
                
                # Show add files interface
                if st.session_state.get('show_add_files', False):
                    st.markdown('<div style="margin-top: 2rem;"></div>', unsafe_allow_html=True)
                    st.markdown("""
                    <div class="glass-panel" style="background: rgba(168, 85, 247, 0.1); border-color: rgba(168, 85, 247, 0.3);">
                        <h4 style="color: #a855f7; margin-bottom: 0.5rem;">➕ Add Files to Existing Database</h4>
                        <p style="margin: 0.5rem 0; opacity: 0.9;">
                            Upload new files to add to the selected database
                        </p>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    add_files = st.file_uploader(
                        "Choose files to add",
                        type=['pdf', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt', 'txt', 'md', 'markdown',
                              'png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif',
                              'mp3', 'wav', 'm4a', 'flac', 'ogg',
                              'mp4', 'avi', 'mov', 'mkv', 'flv'],
                        accept_multiple_files=True,
                        key="add_files_uploader"
                    )
                    
                    if add_files:
                        st.markdown(f"**{len(add_files)} file(s) selected**")
                        
                        if st.button("🚀 Add Files to Database", key="add_files_process"):
                            # First load the database if not loaded
                            if not st.session_state.database_loaded:
                                with st.spinner("Loading database first..."):
                                    db_info = st.session_state.scanned_databases[selected_db]
                                    result = load_existing_database(db_info['path'], db_info['db_folder'])
                                    if result:
                                        vector_store, bm25_retriever, documents, embeddings, metadata = result
                                        st.session_state.vector_store = vector_store
                                        st.session_state.bm25_retriever = bm25_retriever
                                        st.session_state.documents = documents
                                        st.session_state.embeddings = embeddings
                                        st.session_state.metadata = metadata
                                        st.session_state.current_db_info = db_info
                            
                            # Save uploaded files temporarily and add to database
                            with st.spinner("Processing files..."):
                                temp_dir = tempfile.mkdtemp()
                                file_paths = []
                                
                                try:
                                    for add_file in add_files:
                                        temp_path = os.path.join(temp_dir, add_file.name)
                                        with open(temp_path, 'wb') as f:
                                            f.write(add_file.getbuffer())
                                        file_paths.append(Path(temp_path))
                                    
                                    db_info = st.session_state.current_db_info
                                    result = add_files_to_database(
                                        file_paths,
                                        db_info['path'],
                                        db_info['db_folder'],
                                        st.session_state.vector_store,
                                        st.session_state.documents,
                                        st.session_state.embeddings
                                    )
                                    
                                    if result:
                                        vector_store, bm25_retriever, documents, embeddings, metadata = result
                                        st.session_state.vector_store = vector_store
                                        st.session_state.bm25_retriever = bm25_retriever
                                        st.session_state.documents = documents
                                        st.session_state.metadata = metadata
                                        st.success("✅ Files added successfully!")
                                        st.session_state.show_add_files = False
                                        st.balloons()
                                        st.rerun()
                                    
                                except Exception as e:
                                    st.error(f"❌ Error adding files: {e}")

        else:  # Create New Database
            st.markdown('<div style="margin-top: 2rem;"></div>', unsafe_allow_html=True)

            # Initialize session state for files
            if 'found_pdfs' not in st.session_state:
                st.session_state.found_pdfs = []
            if 'uploaded_files' not in st.session_state:
                st.session_state.uploaded_files = []

            st.markdown("""
            <div class="glass-panel">
                <h4 style="margin-bottom: 1rem;">📂 Select Document Source</h4>
                <p style="margin-bottom: 1rem; opacity: 0.8;">
                    Choose how to provide your documents (PDF, Word, Excel, Images, Audio, Video, etc.)
                </p>
            </div>
            """, unsafe_allow_html=True)

            # Source selection
            source_option = st.radio(
                "📁 Document Source",
                ["📤 Upload Files (Drag & Drop)", "📂 Use Server Folder"],
                horizontal=True,
                help="Choose to upload files from your computer or use files from server"
            )

            if source_option == "📤 Upload Files (Drag & Drop)":
                st.markdown('<div style="margin-top: 1.5rem;"></div>', unsafe_allow_html=True)
                
                st.markdown("""
                <div class="glass-panel" style="background: rgba(99, 102, 241, 0.1); border-color: rgba(99, 102, 241, 0.3);">
                    <h4 style="color: #6366f1; margin-bottom: 0.5rem;">📤 Upload Your Documents</h4>
                    <p style="margin: 0.5rem 0; opacity: 0.9;">
                        Drag and drop files below or click to browse
                    </p>
                    <p style="margin: 0; font-size: 0.9rem; opacity: 0.75;">
                        Supported: PDF, Word, Excel, PowerPoint, Images, Audio, Video, Text, Markdown
                    </p>
                </div>
                """, unsafe_allow_html=True)

                uploaded_files = st.file_uploader(
                    "Choose files",
                    type=['pdf', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt', 'txt', 'md', 'markdown',
                          'png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif',
                          'mp3', 'wav', 'm4a', 'flac', 'ogg',
                          'mp4', 'avi', 'mov', 'mkv', 'flv'],
                    accept_multiple_files=True,
                    key="file_uploader",
                    help="Select one or more files from your computer"
                )

                save_folder = st.text_input(
                    "💾 Database Save Location",
                    value="vector_databases",
                    help="Where to save the vector database"
                )

                if uploaded_files:
                    st.markdown('<div style="margin-top: 1.5rem;"></div>', unsafe_allow_html=True)
                    
                    # Show uploaded files preview
                    with st.expander(f"📄 Uploaded Files ({len(uploaded_files)})", expanded=True):
                        for i, uploaded_file in enumerate(uploaded_files, 1):
                            file_size = uploaded_file.size / 1024  # KB
                            st.markdown(f"**{i}.** `{uploaded_file.name}` ({file_size:.1f} KB)")

                    # Process uploaded files
                    if st.button("🔄 Process Uploaded Files", key="process_uploads", width="stretch"):
                        with st.spinner("Saving uploaded files..."):
                            import tempfile
                            temp_dir = tempfile.mkdtemp()
                            pdf_paths = []
                            
                            try:
                                for uploaded_file in uploaded_files:
                                    # Save uploaded file to temporary location
                                    temp_path = os.path.join(temp_dir, uploaded_file.name)
                                    with open(temp_path, 'wb') as f:
                                        f.write(uploaded_file.getbuffer())
                                    pdf_paths.append(Path(temp_path))
                                
                                st.session_state.found_pdfs = pdf_paths
                                st.session_state.pdf_folder = temp_dir
                                st.session_state.save_folder = save_folder
                                st.session_state.uploaded_files = uploaded_files
                                
                                st.success(f"✅ {len(pdf_paths)} file(s) ready for processing!")
                                
                            except Exception as e:
                                st.error(f"❌ Error processing files: {e}")

            else:  # Use Server Folder
                st.markdown('<div style="margin-top: 1.5rem;"></div>', unsafe_allow_html=True)

                # Get current directory and list available folders
                current_dir = os.getcwd()
                
                # Get list of directories in workspace
                available_folders = ["Data"]  # Default folder
                try:
                    for item in os.listdir(current_dir):
                        full_path = os.path.join(current_dir, item)
                        if os.path.isdir(full_path) and not item.startswith('.') and item not in ['vector_databases', 'vector_databases_azhar', '__pycache__', 'assets']:
                            if item not in available_folders:
                                available_folders.append(item)
                except:
                    pass

                col1, col2 = st.columns([2, 1])
                
                with col1:
                    # Folder selection dropdown
                    selected_folder = st.selectbox(
                        "📁 Select Folder",
                        available_folders,
                        help="Choose a folder from the server"
                    )
                    
                    # Option to enter custom path
                    use_custom = st.checkbox("Use custom folder path", key="use_custom_path")
                    
                    if use_custom:
                        pdf_folder = st.text_input(
                            "Custom Folder Path",
                            value=selected_folder,
                            placeholder="/path/to/your/pdfs",
                            help="Enter the full path to your PDF folder"
                        )
                    else:
                        pdf_folder = selected_folder

                with col2:
                    save_folder = st.text_input(
                        "💾 Database Save Location",
                        value="vector_databases",
                        help="Where to save the vector database"
                    )

                st.markdown('<div style="margin-top: 1.5rem;"></div>', unsafe_allow_html=True)

                if st.button("🔍 Scan Selected Folder", key="check_pdfs", width="stretch"):
                    if os.path.exists(pdf_folder):
                        with st.spinner(f"Scanning {pdf_folder} for supported files..."):
                            all_files = get_all_documents(pdf_folder)
                            st.session_state.found_pdfs = all_files
                            st.session_state.pdf_folder = pdf_folder
                            st.session_state.save_folder = save_folder
                            
                            if all_files:
                                # Count file types
                                file_type_counts = {}
                                for f in all_files:
                                    ext = f.suffix.lower()
                                    file_type_counts[ext] = file_type_counts.get(ext, 0) + 1
                                
                                type_summary = ', '.join([f"{count} {ext}" for ext, count in sorted(file_type_counts.items())])
                                st.success(f"✅ Found {len(all_files)} file(s) in {pdf_folder}")
                                st.info(f"📊 File types: {type_summary}")
                                
                                # Show preview of found files
                                with st.expander("📄 View Found Files", expanded=True):
                                    for i, file in enumerate(all_files[:10], 1):  # Show first 10
                                        st.markdown(f"**{i}.** `{file.name}` ({file.suffix}) - {file.parent}")
                                    if len(all_files) > 10:
                                        st.markdown(f"*... and {len(all_files) - 10} more files*")
                            else:
                                st.error(f"❌ No supported files found in {pdf_folder}")
                                st.info("💡 Supported formats: PDF, Word, Excel, PowerPoint, Images, Audio, Video, Text, Markdown")
                    else:
                        st.error(f"❌ Folder does not exist: {pdf_folder}")
                        st.info("💡 Please check the folder path and try again")

            # Show create button if PDFs are found
            if st.session_state.found_pdfs:
                st.markdown('<div style="margin-top: 1.5rem;"></div>', unsafe_allow_html=True)
                
                file_type_counts = {}
                for f in st.session_state.found_pdfs:
                    ext = f.suffix.lower()
                    file_type_counts[ext] = file_type_counts.get(ext, 0) + 1
                type_summary = ', '.join([f"{count} {ext}" for ext, count in sorted(file_type_counts.items())])
                
                st.markdown(f"""
                <div class="glass-panel" style="background: rgba(16, 185, 129, 0.1); border-color: rgba(16, 185, 129, 0.3);">
                    <h4 style="color: #10b981; margin-bottom: 0.5rem;">✅ Ready to Process</h4>
                    <p style="margin: 0; opacity: 0.9;">
                        <strong>{len(st.session_state.found_pdfs)} files</strong> will be processed from <strong>{st.session_state.pdf_folder}</strong>
                    </p>
                    <p style="margin: 0.5rem 0 0 0; font-size: 0.9rem; opacity: 0.75;">
                        File types: {type_summary}
                    </p>
                </div>
                """, unsafe_allow_html=True)

                if st.button("🚀 Create Vector Database", key="create_db", width="stretch"):
                    result = create_new_database(
                        st.session_state.found_pdfs,
                        st.session_state.save_folder
                    )

                    if result:
                        vector_store, bm25_retriever, documents, embeddings, metadata = result
                        st.session_state.vector_store = vector_store
                        st.session_state.bm25_retriever = bm25_retriever
                        st.session_state.documents = documents
                        st.session_state.embeddings = embeddings
                        st.session_state.metadata = metadata
                        st.session_state.database_loaded = True
                        st.session_state.found_pdfs = []  # Reset
                        st.success("✅ Database created successfully! Switch to AI Model tab to configure.")
                        st.balloons()

    with config_tab2:
        st.markdown('<div class="card-title"><span class="card-title-icon">🤖</span>AI Model Configuration</div>', unsafe_allow_html=True)

        if not st.session_state.database_loaded:
            st.warning("⚠️ Please load or create a database first")
        else:
            st.markdown('<div style="margin-top: 2rem;"></div>', unsafe_allow_html=True)

            model_provider = st.selectbox(
                "🏭 Model Provider",
                ["Local Models (Ollama)", "API Models (Groq)"]
            )

            if model_provider == "Local Models (Ollama)":
                ollama_running = check_ollama_connection()
                if ollama_running:
                    available_models = get_available_ollama_models()
                    if available_models:
                        st.success("✅ Ollama is running")
                        selected_model = st.selectbox("Select Model", available_models)
                    else:
                        st.error("❌ No models found. Run: `ollama pull llama3.2:3b`")
                        selected_model = None
                else:
                    st.error("❌ Ollama is not running")
                    selected_model = None
            else:
                if groq_api_key:
                    st.success("✅ Groq API Key found")
                    selected_model_display = st.selectbox(
                        "Select Model",
                        list(AVAILABLE_MODELS["API Models (Groq)"].keys())
                    )
                    selected_model = AVAILABLE_MODELS["API Models (Groq)"][selected_model_display]
                else:
                    st.error("❌ No API key found")
                    selected_model = None

            if selected_model and st.button("🚀 Initialize AI Model", key="init_model"):
                with st.spinner("Initializing..."):
                    llm = initialize_llm(model_provider, selected_model, groq_api_key)

                    if llm:
                        st.session_state.llm = llm
                        st.session_state.model_provider = model_provider
                        st.session_state.selected_model = selected_model
                        st.session_state.llm_configured = True
                        st.success("✅ AI Model initialized successfully! Switch to Chat tab to start.")
                        st.balloons()

# TAB 3: CHAT
with tab3:
    if not st.session_state.database_loaded or not st.session_state.llm_configured:
        st.markdown("""
        <div class="modern-card">
            <h3 style="text-align: center;">⚠️ Configuration Required</h3>
            <p style="text-align: center; margin-top: 1rem; opacity: 0.8;">
                Please configure your database and AI model in the Configuration tab before starting a chat.
            </p>
        </div>
        """, unsafe_allow_html=True)
    else:
        st.markdown('<div class="card-title"><span class="card-title-icon">💬</span>Intelligent Conversation</div>', unsafe_allow_html=True)

        # Create hybrid retriever
        hybrid_retriever = create_hybrid_retriever(
            st.session_state.vector_store,
            st.session_state.bm25_retriever
        )

        llm = st.session_state.llm

        prompt = ChatPromptTemplate.from_template("""
        You are an expert assistant analyzing documents. Answer the following question based only on the provided context.
        Be comprehensive, detailed, and include specific references.

        <context>
        {context}
        </context>

        Question: {input}

        Answer:
        """)

        document_chain = create_stuff_documents_chain(llm, prompt)
        retrieval_chain = create_retrieval_chain(hybrid_retriever, document_chain)

        # Display chat history
        for idx, message in enumerate(st.session_state.messages):
            with st.chat_message(message["role"]):
                if message["role"] == "assistant":
                    st.markdown(message["content"], unsafe_allow_html=True)

                    # Display sources if available
                    if "sources" in message and message["sources"]:
                        with st.expander("📚 View Sources", expanded=False):
                            import pandas as pd
                            # Create simplified dataframe for display
                            display_sources = []
                            for s in message["sources"]:
                                display_sources.append({
                                    "Citation": s.get("Citation", ""),
                                    "Document": s.get("Document", "Unknown"),
                                    "Page": str(s.get("Page", "Unknown"))
                                })
                            df = pd.DataFrame(display_sources)
                            st.dataframe(df, use_container_width=True, hide_index=True)

                            # Show detailed source content
                            st.markdown("---")
                            st.markdown("**📖 Source Content Preview:**")
                            for i, source in enumerate(message["sources"]):
                                st.markdown(f"**[{i+1}] {source.get('Document', 'Unknown')} - Page {source.get('Page', 'Unknown')}**")
                                if "content" in source:
                                    preview = source["content"][:300] + "..." if len(source["content"]) > 300 else source["content"]
                                    st.markdown(f'<div style="background: rgba(255,255,255,0.03); padding: 1rem; border-radius: 8px; margin: 0.5rem 0 1rem 0; color: rgba(255,255,255,0.7); line-height: 1.6;">{preview}</div>', unsafe_allow_html=True)
                else:
                    st.markdown(message["content"])

        # Chat input at the bottom
        if prompt_input := st.chat_input("💭 Ask me anything about your documents..."):
            # Add user message
            st.session_state.messages.append({"role": "user", "content": prompt_input})

            # Generate response
            try:
                response = retrieval_chain.invoke({"input": prompt_input})
                answer = response["answer"]

                # Process answer with citations
                processed_answer = process_answer_with_citations(answer, response["context"])
                styled_answer = f'<div class="answer-content">{processed_answer}</div>'

                # Prepare source data
                import pandas as pd
                source_data = []
                for i, doc in enumerate(response["context"]):
                    source_file = doc.metadata.get('source_file', 'Unknown')
                    page_num = doc.metadata.get('page', 'Unknown')
                    full_path = doc.metadata.get('full_path', '')
                    source_data.append({
                        "Citation": f"[{i+1}]",
                        "Document": source_file,
                        "Page": str(page_num),
                        "content": doc.page_content,
                        "path": full_path
                    })

                # Save to session with sources
                st.session_state.messages.append({
                    "role": "assistant",
                    "content": styled_answer,
                    "sources": source_data
                })

                # Force rerun to display the new message
                st.rerun()

            except Exception as e:
                error_msg = f"❌ Error: {str(e)}"
                st.error(error_msg)
                st.session_state.messages.append({"role": "assistant", "content": error_msg})
                st.rerun()

        # Action buttons
        st.markdown('<div style="margin-top: 2rem;"></div>', unsafe_allow_html=True)
        col1, col2, col3 = st.columns([1, 1, 1])
        with col1:
            if st.button("🗑️ Clear Chat", key="clear_chat"):
                st.session_state.messages = []
                st.rerun()
        with col2:
            if st.button("🔄 Reset All", key="reset_config"):
                st.session_state.database_loaded = False
                st.session_state.llm_configured = False
                st.session_state.messages = []
                st.session_state.scanned_databases = []
                st.session_state.found_pdfs = []
                st.rerun()

st.markdown('</div>', unsafe_allow_html=True)
